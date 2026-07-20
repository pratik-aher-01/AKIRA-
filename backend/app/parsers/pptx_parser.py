from pptx import Presentation

def extract_text_from_pptx(file_path: str) -> str:
    """
    Reads a PowerPoint file and extracts text from slides and speaker notes.
    """
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            # 1. Get text from shapes (text boxes, titles) on the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            
            # 2. Get text from speaker notes if they exist
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                text += slide.notes_slide.notes_text_frame.text + "\n"
                
            text += "\n" # Add a space between slides
        return text.strip()
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
        return ""